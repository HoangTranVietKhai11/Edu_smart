import os
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from typing import List, Dict, Any
from dotenv import load_dotenv

load_dotenv()

CHROMA_PERSIST_DIR = os.getenv("CHROMA_PERSIST_DIR", "./chroma_db")
EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL", "all-MiniLM-L6-v2")
CHUNK_SIZE = 500
CHUNK_OVERLAP = 50

# Initialize ChromaDB
chroma_client = chromadb.PersistentClient(path=CHROMA_PERSIST_DIR)
collection = chroma_client.get_or_create_collection(
    name="edusmart_documents",
    metadata={"hnsw:space": "cosine"},
)

# Initialize embedding model
print(f"Loading embedding model: {EMBEDDING_MODEL}")
embedding_model = SentenceTransformer(EMBEDDING_MODEL)


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF using PyMuPDF."""
    text = ""
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
        doc.close()
    except Exception as e:
        print(f"PDF extraction error: {e}")
    return text.strip()


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX."""
    text = ""
    try:
        doc = DocxDocument(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"DOCX extraction error: {e}")
    return text.strip()


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"PPTX extraction error: {e}")
    return text.strip()


def extract_text(file_path: str, file_type: str) -> str:
    """Extract text based on file type."""
    extractors = {
        "pdf": extract_text_from_pdf,
        "docx": extract_text_from_docx,
        "pptx": extract_text_from_pptx,
    }
    extractor = extractors.get(file_type)
    if not extractor:
        return ""
    return extractor(file_path)


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Split text into overlapping chunks."""
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
    return chunks


def embed_document(document_id: str, file_path: str, file_type: str, metadata: Dict) -> List[str]:
    """
    Full RAG pipeline: Extract → Chunk → Embed → Store in ChromaDB.
    Returns list of vector IDs.
    """
    # Extract text
    text = extract_text(file_path, file_type)
    if not text:
        raise ValueError(f"Could not extract text from {file_path}")

    # Chunk text
    chunks = chunk_text(text)
    if not chunks:
        raise ValueError("No text chunks generated")

    # Generate embeddings
    embeddings = embedding_model.encode(chunks, show_progress_bar=False).tolist()

    # Store in ChromaDB
    vector_ids = [f"{document_id}_chunk_{i}" for i in range(len(chunks))]
    chunk_metadata = [
        {
            **metadata,
            "document_id": document_id,
            "chunk_index": i,
            "chunk_total": len(chunks),
        }
        for i in range(len(chunks))
    ]

    collection.add(
        ids=vector_ids,
        embeddings=embeddings,
        documents=chunks,
        metadatas=chunk_metadata,
    )

    print(f"✅ Embedded {len(chunks)} chunks for document {document_id}")
    return vector_ids


def retrieve_context(query: str, class_id: str = None, n_results: int = 5) -> List[Dict]:
    """
    Semantic retrieval from ChromaDB.
    Returns list of relevant chunks with metadata.
    """
    query_embedding = embedding_model.encode([query]).tolist()

    where_filter = {}
    if class_id:
        where_filter["class_id"] = class_id

    results = collection.query(
        query_embeddings=query_embedding,
        n_results=n_results,
        where=where_filter if where_filter else None,
    )

    context_items = []
    if results and results["documents"]:
        for i, doc in enumerate(results["documents"][0]):
            context_items.append({
                "content": doc,
                "metadata": results["metadatas"][0][i] if results["metadatas"] else {},
                "distance": results["distances"][0][i] if results["distances"] else 1.0,
            })

    return context_items


def delete_document_vectors(document_id: str):
    """Remove all vectors for a document."""
    try:
        all_ids = collection.get(where={"document_id": document_id})["ids"]
        if all_ids:
            collection.delete(ids=all_ids)
        print(f"Deleted vectors for document {document_id}")
    except Exception as e:
        print(f"Error deleting vectors: {e}")
